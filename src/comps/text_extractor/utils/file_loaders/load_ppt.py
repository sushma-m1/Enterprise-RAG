# Copyright (C) 2024-2025 Intel Corporation
# SPDX-License-Identifier: Apache-2.0

import olefile
import os
import pptx
import uuid

from pptx.exc import PackageNotFoundError

from comps.text_extractor.utils.file_loaders.abstract_loader import AbstractLoader
from comps.text_extractor.utils.file_loaders.load_image import LoadImage
from comps.cores.mega.logger import get_opea_logger

logger = get_opea_logger(f"{__file__.split('comps/')[1].split('/', 1)[0]}_microservice")

class LoadPpt(AbstractLoader):
    """ Load and parse ppt/pptx files. """
    def is_encrypted(self):
        """ Check if the ppt file is encrypted. """
        ole = olefile.OleFileIO(self.file_path)
        if ole.exists('EncryptedPackage'):
            logger.info(f"[{self.file_path}] This file is encrypted.")
            return True
        return False

    def extract_text(self):
        converted = False
        if self.file_type == "ppt":
            self.convert_to_pptx(self.file_path)
            converted = True

        try:
            text = ""
            prs = pptx.Presentation(self.file_path)
            logger.info(f"[{self.file_path}] Processing {len(prs.slides)} slides")

            for slide in prs.slides:
                for shape in sorted(slide.shapes, key=lambda shape: (shape.top, shape.left)):
                    if shape.has_text_frame:
                        if shape.text:
                            text += shape.text + "\n"
                    if shape.has_table:
                        table_contents = "\n".join(
                            [
                                "\t".join([(cell.text if hasattr(cell, "text") else "") for cell in row.cells])
                                for row in shape.table.rows
                                if hasattr(row, "cells")
                            ]
                        )
                        if table_contents:
                            text += table_contents + "\n"
                    if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                        img_path = ""
                        try:
                            img_path = self.save_image(shape.image, os.getenv("UPLOAD_PATH", "/tmp/opea_upload"))
                            logger.debug(f"[{self.file_path}] Extracted {img_path} for processing")
                            img_loader = LoadImage(img_path)
                            text += img_loader.extract_text()
                            logger.info(f"[{self.file_path}] Processed image {img_path}")
                        except Exception as e:
                            logger.error(f"[{self.file_path}] Error parsing image: {e}. Ignoring...")
                        finally:
                            if img_path and img_path != "" and os.path.exists(img_path) and not os.path.isdir(img_path):
                                logger.debug(f"[{self.file_path}] Removed {img_path} after processing")
                                os.remove(img_path)
        except PackageNotFoundError as e:
            if os.path.exists(self.file_path) and self.is_encrypted():
                logger.error(f"[{self.file_path}] The file is encrypted and cannot be processed: {e}")
                raise ValueError(f"The file {self.file_path} is encrypted and cannot be processed. Please decrypt it before uploading.")
            else:
                logger.error(f"[{self.file_path}] Error opening PPTX file: {e}")
                raise e
        finally:
            # Ensure the converted file is deleted. The original file deletion
            # should be handled in the caller code
            if converted and os.path.exists(self.file_path):
                os.remove(self.file_path)
                logger.info(f"Removed temporary converted file {self.file_path}")

        return text

    def save_image(self, image, save_path="/tmp/opea_upload"):
        import uuid
        import os
        if image:
            image_filename = os.path.join(save_path, f"{uuid.uuid4()}.{image.ext}")
            with open(image_filename, 'wb') as f:
                f.write(image.blob)
            return image_filename
        return None

    def convert_to_pptx(self, ppt_path):
        """Convert ppt file to pptx file."""
        pptx_path = ppt_path + "x"
        convert_log_file = f'/tmp/convert_{uuid.uuid4()}.log'
        exit_code = os.system(f"libreoffice --headless --invisible --convert-to pptx --outdir {os.path.dirname(pptx_path)} '{ppt_path}' > {convert_log_file} 2>&1")
        if exit_code != 0 or not os.path.exists(pptx_path):
            error = ""
            logger.error(f"Failed to convert {ppt_path} to pptx format. Exit code: {exit_code}")
            if os.path.exists(convert_log_file):
                try:
                    with open(convert_log_file, 'r') as f:
                        error = f.read()
                    logger.error(f"Conversion error: {error}")
                finally:
                    os.remove(convert_log_file)
            raise ValueError(f"Failed to convert {ppt_path} to pptx format. Error: {error}")
        else:
            self.file_path = pptx_path
            self.file_type = "pptx"
