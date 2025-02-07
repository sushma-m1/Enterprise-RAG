# Copyright (C) 2024-2025 Intel Corporation
# SPDX-License-Identifier: Apache-2.0

import os
import pptx
from comps.dataprep.utils.file_loaders.abstract_loader import AbstractLoader


class LoadPpt(AbstractLoader):
    """ Load and parse ppt/pptx files. """
    def extract_text(self):
        if self.file_type == "ppt":
            self.convert_to_pptx(self.file_path)

        text = ""
        prs = pptx.Presentation(self.file_path)
        for slide in prs.slides:
            for shape in sorted(slide.shapes, key=lambda shape: (shape.top, shape.left)):
                if shape.has_text_frame:
                    if shape.text:
                        text += shape.text + "\n"
                if shape.has_table:
                    table_contents = "\n".join(
                        [
                            "\t".join([(cell.text if hasattr(cell, "text") else "") for cell in row.cells])
                            for row in shape.table.rows
                            if hasattr(row, "cells")
                        ]
                    )
                    if table_contents:
                        text += table_contents + "\n"

        return text

    def convert_to_pptx(self, ppt_path):
        """Convert ppt file to pptx file."""
        pptx_path = ppt_path + "x"
        os.system(f"libreoffice --headless --invisible --convert-to pptx --outdir {os.path.dirname(pptx_path)} {ppt_path}")
        self.file_path = pptx_path
        self.file_type = "pptx"
